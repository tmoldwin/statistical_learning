"""Build Google Slides–compatible decks from paper/draft.md figures.

Outputs under ``paper/slides/``:
- ``paper_figures.pptx`` — figure slides only (optional generated title)
- ``paper_deck.pptx`` — ``intro_front.pptx`` (narrative front) + figure slides

Import in Google Slides: File → Import slides → Upload.
"""

from __future__ import annotations

import argparse
import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

REPO_ROOT = Path(__file__).resolve().parents[1]
DRAFT = REPO_ROOT / "paper" / "draft.md"
OUT_DIR = REPO_ROOT / "paper" / "slides"
OUT_PPTX = OUT_DIR / "paper_figures.pptx"
OUT_DECK = OUT_DIR / "paper_deck.pptx"
INTRO_FRONT = OUT_DIR / "intro_front.pptx"

IMG_RE = re.compile(r"!\[([^\]]*)\]\(([^)]+)\)")
CAPTION_RE = re.compile(r"^\*\*Figure (\d+)\.\*\*\s*(.+)$", re.MULTILINE)
TITLE_RE = re.compile(r"^#\s+(.+)$", re.MULTILINE)


def parse_figures(draft_text: str) -> list[dict]:
    """Pair each draft image with figure number and alt-text title."""
    lines = draft_text.splitlines()
    figures: list[dict] = []
    i = 0
    while i < len(lines):
        m_img = IMG_RE.match(lines[i].strip())
        if not m_img:
            i += 1
            continue
        alt, rel_path = m_img.group(1), m_img.group(2)
        fig_num: int | None = None
        j = i + 1
        while j < len(lines) and not lines[j].strip():
            j += 1
        if j < len(lines):
            m_cap = CAPTION_RE.match(lines[j].strip())
            if m_cap:
                fig_num = int(m_cap.group(1))
        figures.append({
            "num": fig_num if fig_num is not None else len(figures) + 1,
            "title": alt.strip() or f"Figure {fig_num or len(figures) + 1}",
            "path": (REPO_ROOT / "paper" / rel_path).resolve(),
        })
        i = j + 1 if j > i else i + 1
    return figures


def _layout_by_name(prs: Presentation, *names: str):
    wanted = {n.upper() for n in names}
    for layout in prs.slide_layouts:
        if layout.name.upper() in wanted:
            return layout
    return prs.slide_layouts[min(6, len(prs.slide_layouts) - 1)]


def _add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    # Default blank template: layout 0 is the title slide. Named lookup is
    # unreliable across Google-exported masters (TITLE may lack a title shape).
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    if slide.shapes.title is not None:
        slide.shapes.title.text = title
    else:
        box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.0), float(prs.slide_width) - Inches(1.0), Inches(1.2),
        )
        p = box.text_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(28)
        p.font.bold = True
    if len(slide.placeholders) > 1:
        try:
            slide.placeholders[1].text = subtitle
        except Exception:
            pass


def _fit_image_box(img_w: int, img_h: int, max_w: float, max_h: float) -> tuple[float, float]:
    if img_w <= 0 or img_h <= 0:
        return max_w, max_h
    scale = min(max_w / img_w, max_h / img_h)
    return img_w * scale, img_h * scale


def _add_figure_slide(
    prs: Presentation,
    *,
    title: str,
    image_path: Path,
) -> None:
    blank = _layout_by_name(prs, "BLANK")
    slide = prs.slides.add_slide(blank)

    slide_w = float(prs.slide_width)
    slide_h = float(prs.slide_height)
    # Scale chrome to whatever slide size the deck uses (intro is 10×5.625).
    margin_x = slide_w * 0.02
    top_y = slide_h * 0.03
    title_h = slide_h * 0.12
    bottom_margin = slide_h * 0.03
    title_pt = 18 if slide_h < Inches(6.5) else 24

    title_box = slide.shapes.add_textbox(margin_x, top_y, slide_w - 2 * margin_x, title_h)
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(title_pt)
    p.font.bold = True

    img_top = top_y + title_h + slide_h * 0.01
    img_max_w = slide_w - 2 * margin_x
    img_max_h = slide_h - img_top - bottom_margin

    from PIL import Image

    with Image.open(image_path) as im:
        px_w, px_h = im.size
    disp_w, disp_h = _fit_image_box(px_w, px_h, img_max_w, img_max_h)
    left = (slide_w - disp_w) / 2
    top = img_top + (img_max_h - disp_h) / 2
    slide.shapes.add_picture(str(image_path), left, top, width=disp_w, height=disp_h)


def _append_figures(prs: Presentation, figures: list[dict]) -> list[str]:
    missing: list[str] = []
    for fig in figures:
        path = fig["path"]
        if not path.is_file():
            missing.append(str(path.relative_to(REPO_ROOT)))
            continue
        _add_figure_slide(prs, title=str(fig["title"]), image_path=path)
    return missing


def build_slides(
    figures: list[dict],
    *,
    title: str,
    subtitle: str,
    slide_width_in: float = 13.333,
    slide_height_in: float = 7.5,
    include_title: bool = True,
) -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(slide_width_in)
    prs.slide_height = Inches(slide_height_in)

    if include_title:
        _add_title_slide(prs, title, subtitle)

    missing = _append_figures(prs, figures)
    if missing:
        print("missing images (skipped):")
        for m in missing:
            print(f"  {m}")
    return prs


def build_deck_from_intro(
    figures: list[dict],
    *,
    intro_path: Path,
) -> Presentation:
    """Open the narrative front deck and append figure slides (no extra title)."""
    if not intro_path.is_file():
        raise FileNotFoundError(intro_path)
    # Caller must save to a different path so the intro source stays intact.
    prs = Presentation(str(intro_path))
    missing = _append_figures(prs, figures)
    if missing:
        print("missing images (skipped):")
        for m in missing:
            print(f"  {m}")
    return prs


def main(argv: list[str] | None = None) -> None:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument(
        "--intro",
        type=Path,
        default=INTRO_FRONT,
        help=f"Narrative front PPTX to prepend (default: {INTRO_FRONT})",
    )
    parser.add_argument(
        "--no-deck",
        action="store_true",
        help="Only write paper_figures.pptx (skip combined paper_deck.pptx)",
    )
    args = parser.parse_args(argv)

    draft_text = DRAFT.read_text(encoding="utf-8")
    title_m = TITLE_RE.search(draft_text)
    paper_title = title_m.group(1) if title_m else "Paper figures"
    subtitle = "Generated from paper/draft.md · import into Google Slides via Upload"

    figures = parse_figures(draft_text)
    if not figures:
        print(f"no figures found in {DRAFT}", file=sys.stderr)
        sys.exit(1)

    OUT_DIR.mkdir(parents=True, exist_ok=True)

    # Figures-only deck: match intro slide size when available so imports align.
    w_in, h_in = 13.333, 7.5
    if args.intro.is_file():
        intro_prs = Presentation(str(args.intro))
        w_in = float(intro_prs.slide_width.inches)
        h_in = float(intro_prs.slide_height.inches)

    def _save(prs: Presentation, path: Path) -> Path:
        try:
            prs.save(str(path))
            return path
        except PermissionError:
            alt = path.with_name(path.stem + "_new" + path.suffix)
            prs.save(str(alt))
            print(f"warning: {path.name} locked; wrote {alt.name} instead", flush=True)
            return alt

    # Combined deck first (primary deliverable for Google Slides import).
    if not args.no_deck and args.intro.is_file():
        n_front = len(Presentation(str(args.intro)).slides)
        deck = build_deck_from_intro(figures, intro_path=args.intro)
        deck_path = _save(deck, OUT_DECK)
        n_total = len(deck.slides)
        print(
            f"wrote {deck_path} "
            f"({n_front} intro + {n_total - n_front} figures = {n_total} slides)",
            flush=True,
        )
    elif not args.no_deck:
        print(f"no intro deck at {args.intro}; skipped {OUT_DECK.name}", flush=True)

    figs_prs = build_slides(
        figures,
        title=paper_title,
        subtitle=subtitle,
        slide_width_in=w_in,
        slide_height_in=h_in,
        include_title=True,
    )
    figs_path = _save(figs_prs, OUT_PPTX)
    print(f"wrote {figs_path} ({len(figs_prs.slides) - 1} figure slides + title)")


if __name__ == "__main__":
    main()
